#!/usr/bin/env python3
"""
Generador de presentación PowerPoint
Crea automáticamente una presentación PPTX con el contenido de Marketing 1.0-6.0
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    print("Instalando python-pptx...")
    import subprocess
    subprocess.check_call(["pip", "install", "-q", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colores por versión (Pastel para accesibilidad)
COLORS = {
    1: RGBColor(240, 212, 184),    # Naranja pastel
    2: RGBColor(184, 224, 240),    # Azul pastel
    3: RGBColor(240, 184, 176),    # Rojo pastel
    4: RGBColor(216, 186, 232),    # Púrpura pastel
    5: RGBColor(168, 232, 216),    # Turquesa pastel
    6: RGBColor(184, 232, 176),    # Verde pastel
}

def add_title_slide(prs, title, subtitle):
    """Agrega diapositiva de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(200, 186, 232)
    
    # Título
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, version, title, period, color):
    """Agrega diapositiva de sección para cada Marketing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Barra de color
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    
    # Versión
    version_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(2), Inches(0.6))
    version_frame = version_box.text_frame
    p = version_frame.paragraphs[0]
    p.text = f"Marketing {version}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Período
    period_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
    period_frame = period_box.text_frame
    p = period_frame.paragraphs[0]
    p.text = period
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)

def add_content_slide(prs, title, content, color):
    """Agrega diapositiva con contenido"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Encabezado
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.color.rgb = color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Contenido
    left = Inches(0.7)
    top = Inches(1.0)
    width = Inches(8.6)
    height = Inches(5.8)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for line in content:
        if line.startswith('BOLD:'):
            p = text_frame.add_paragraph()
            p.text = line.replace('BOLD:', '')
            p.font.bold = True
            p.font.size = Pt(18)
            p.space_before = Pt(8)
        else:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.level = 0
            p.space_after = Pt(4)

def add_chart_slide(prs, title, color):
    """Agrega diapositiva con gráfico de evolución"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Encabezado
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.color.rgb = color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Datos del gráfico
    chart_data = ChartData()
    chart_data.categories = ['1900', '1950', '2000', '2010', '2015', '2020', '2026']
    chart_data.add_series('Evolución del Marketing', (1, 2, 3, 4, 5, 6, 6))
    
    # Agregar gráfico
    x, y, cx, cy = Inches(1), Inches(1.2), Inches(8), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True

def add_timeline_slide(prs, title, color):
    """Agrega diapositiva con timeline de versiones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Encabezado
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.color.rgb = color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Línea de timeline
    line = slide.shapes.add_shape(1, Inches(1), Inches(3.5), Inches(8), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    # Puntos en la timeline
    versions = [
        ("1.0", "1900-1950", COLORS[1]),
        ("2.0", "1950-2000", COLORS[2]),
        ("3.0", "2000-2010", COLORS[3]),
        ("4.0", "2010-2015", COLORS[4]),
        ("5.0", "2015-2020", COLORS[5]),
        ("6.0", "2020-2026", COLORS[6])
    ]
    
    for i, (ver, period, col) in enumerate(versions):
        x_pos = Inches(1 + i * 1.2)
        # Punto
        point = slide.shapes.add_shape(3, x_pos - Inches(0.1), Inches(3.3), Inches(0.2), Inches(0.2))
        point.fill.solid()
        point.fill.fore_color.rgb = col
        
        # Texto versión
        text_box = slide.shapes.add_textbox(x_pos - Inches(0.3), Inches(1.8), Inches(0.6), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = ver
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        
        # Texto período
        period_box = slide.shapes.add_textbox(x_pos - Inches(0.4), Inches(4.5), Inches(0.8), Inches(0.5))
        period_frame = period_box.text_frame
        p = period_frame.paragraphs[0]
        p.text = period
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

# PORTADA
add_title_slide(prs, 
    "Historia y Evolución del Marketing",
    "6 Generaciones de Estrategias Comerciales (1900-2026)")

# GRÁFICO DE EVOLUCIÓN
add_chart_slide(prs, "📈 Evolución del Marketing a través del Tiempo", RGBColor(200, 186, 232))

# TIMELINE DE VERSIONES
add_timeline_slide(prs, "⏰ Timeline de las Generaciones del Marketing", RGBColor(200, 186, 232))

# SECCIÓN: MARKETING VS VENTAS COMPARACIÓN
add_content_slide(prs, "🔄 Marketing vs Ventas: Diferencia Clave", [
    'BOLD:MARKETING:',
    '★ Concepto: Crear, comunicar y entregar valor',
    '★ Horizonte: Semanas a años (estratégico)',
    '★ Métrica: Brand awareness, Lead generation',
    '', 
    'BOLD:VENTAS:',
    '★ Concepto: Convertir oportunidades en compras',
    '★ Horizonte: Días a semanas (táctico)',
    '★ Métrica: % de conversión, Ticket promedio'
], RGBColor(200, 186, 232))

add_content_slide(prs, "Marketing vs Ventas: Clasificación", [
    'BOLD:TIPOS DE MARKETING:',
    '• Digital: SEO, SEM, Social Media, Email',
    '• Tradicional: TV, Radio, Prensa',
    '• Influencia: Influencers, Embajadores',
    '• Contenido: Blogs, Videos, Podcasts',
    '', 
    'BOLD:TIPOS DE VENTAS:',
    '• Directa: B2C, Tiendas, Vendedores',
    '• Corporativa: B2B, Account Managers',
    '• Online: E-commerce, Chatbots',
    '• Automatizada: AI Assistants, CRM'
], RGBColor(200, 186, 232))

add_content_slide(prs, "En Marketing 6.0: La Convergencia", [
    'BOLD:Data + Decisión Smart',
    'Marketing proporciona insights, Ventas cierra con inteligencia',
    '',
    'BOLD:Automatización Inteligente',
    'Bots califican leads, cierran tratos sin intervención',
    '',
    'BOLD:Experiencia Unificada',
    'El cliente vive un flujo continuo awareness → compra → lealtad',
    '',
    'BOLD:Métricas Integradas',
    'ROI se calcula en el journey completo, no por departamento'
], RGBColor(200, 186, 232))

# CONTENIDO DE CADA VERSIÓN
versions_data = [
    {
        'num': 1,
        'title': 'Era de Producción',
        'period': '1900-1950',
        'color': COLORS[1],
        'sections': [
            'Concepto Principal',
            'BOLD:Enfoque en la producción masiva',
            'Bajo poder del consumidor',
            'Publicidad unidireccional',
            '',
            'Características Clave',
            'BOLD:Canales Principales:',
            'Radio, Periódicos, Carteles publicitarios',
            '',
            'Ejemplos: Ford, Coca-Cola, Nestlé'
        ]
    },
    {
        'num': 2,
        'title': 'Era Digital Temprana',
        'period': '1950-2000',
        'color': COLORS[2],
        'sections': [
            'Cambio de enfoque al consumidor',
            'BOLD:Segmentación de mercados',
            'Personalización relativa',
            'Construcción de lealtad',
            '',
            'Canales Principales',
            'BOLD:Televisión, Internet temprano',
            'Email marketing, Tiendas especializadas',
            '',
            'Ejemplos: Amazon, McDonald\'s, Dell'
        ]
    },
    {
        'num': 3,
        'title': 'Era de Valores',
        'period': '2000-2010',
        'color': COLORS[3],
        'sections': [
            'Marketing con propósito y valores',
            'BOLD:Responsabilidad social',
            'Sostenibilidad ambiental',
            'Autenticidad y transparencia',
            '',
            'Canales Principales',
            'BOLD:Redes sociales, Blogs, Marketing viral',
            'Partnerships con ONGs',
            '',
            'Ejemplos: TOMS Shoes, Patagonia, Ben & Jerry\'s'
        ]
    },
    {
        'num': 4,
        'title': 'Era Omnicanal',
        'period': '2010-2015',
        'color': COLORS[4],
        'sections': [
            'Integración online + offline',
            'BOLD:Experiencia del consumidor',
            'Big Data y analítica',
            'Marketing en tiempo real',
            '',
            'Canales Principales',
            'BOLD:E-commerce, Mobile, Social media',
            'Marketing automation, Tiendas conectadas',
            '',
            'Ejemplos: Nike, Starbucks, Sephora'
        ]
    },
    {
        'num': 5,
        'title': 'Era de IA y Automación',
        'period': '2015-2020',
        'color': COLORS[5],
        'sections': [
            'Inteligencia artificial y Machine Learning',
            'BOLD:Personalización 1-a-1 a escala',
            'Automatización inteligente',
            'Marketing predictivo',
            '',
            'Canales Principales',
            'BOLD:Chatbots, Publicidad programática',
            'Recomendación inteligente, Voice commerce',
            '',
            'Ejemplos: Netflix, Amazon Alexa, Google Ads'
        ]
    },
    {
        'num': 6,
        'title': 'Era de Sostenibilidad',
        'period': '2020-2026',
        'color': COLORS[6],
        'sections': [
            'Tecnología + Humanidad + Sostenibilidad',
            'BOLD:Economía circular',
            'Marketing ético y transparente',
            'Impacto social y ambiental',
            '',
            'Canales Principales',
            'BOLD:Comunidades digitales, Metaverso',
            'Live streaming, Blockchain',
            '',
            'Ejemplos: Unilever, Decathlon, Airbnb'
        ]
    }
]

# Agregar secciones
for v in versions_data:
    add_section_slide(prs, f"{v['num']}.0", v['title'], v['period'], v['color'])
    add_content_slide(prs, f"Marketing {v['num']}.0 - {v['title']}", v['sections'], v['color'])

# RESUMEN DE DIFERENCIAS MARKETING VS VENTAS AL FINAL
add_content_slide(prs, "🔄 Resumen: Diferencias Clave entre Marketing y Ventas", [
    'BOLD:MARKETING (Crear Deseo):',
    '• Enfoque: Largo plazo, estratégico',
    '• Objetivo: Generar leads y awareness',
    '• Métricas: ROI de marca, engagement',
    '• Horizonte: Semanas a años',
    '',
    'BOLD:VENTAS (Cerrar Negocio):',
    '• Enfoque: Corto plazo, táctico',
    '• Objetivo: Convertir leads en compras',
    '• Métricas: Conversión, revenue',
    '• Horizonte: Días a semanas',
    '',
    'BOLD:En Marketing 6.0: Convergencia',
    'Ambos departamentos colaboran para una experiencia unificada del cliente'
], RGBColor(200, 186, 232))

# Diapositiva de conclusiones
add_content_slide(prs, "Conclusiones Clave", [
    'BOLD:1. Evolución Constante',
    'El marketing evolucionó desde enfoque en producción a sostenibilidad',
    '',
    'BOLD:2. Tecnología + Humanidad',
    'La mejor estrategia equilibra automatización con conexión genuina',
    '',
    'BOLD:3. El Poder está en el Consumidor',
    'Las empresas ahora deben adaptarse a las demandas del mercado',
    '',
    'BOLD:4. Ética y Sostenibilidad son Competitivas',
    'El futuro del marketing es responsable y transparente'
], COLORS[1])

# Guardar
output_file = 'Historia_y_Evolucion_del_Marketing.pptx'
prs.save(output_file)
print(f"✅ PowerPoint creado: {output_file}")
